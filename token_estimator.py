#!/usr/bin/env python3
"""
Professional Tasks Context Range Estimator.

Estimates source-file text tokens and lower-confidence visual/image tokens,
checks the Professional Tasks minimum reading requirement, and separately
checks an editable maximum context ceiling.

Project defaults:
- Minimum required files: 10
- Minimum required source context: 256,000 tokens
- Maximum selected context: 1,000,000 tokens (provisional/editable)

The 256k minimum is evaluated on required source files only. Prompt tokens do
not count toward that minimum. The maximum check uses the full uploaded source
set plus prompt tokens as a conservative/worst-case estimate.
"""

from __future__ import annotations

import argparse
import csv
import io
import json
import math
import re
import shutil
import tempfile
import xml.etree.ElementTree as ET
import zipfile
import warnings
from dataclasses import asdict, dataclass
from pathlib import Path, PurePosixPath
from functools import lru_cache
from html.parser import HTMLParser
from typing import Callable, Iterable, List, Optional, Sequence, TextIO

from bs4 import BeautifulSoup
from pptx import Presentation


# openpyxl emits noisy warnings for OOXML extensions/styles that this read-only token
# estimator intentionally does not use. Suppress only these known non-fatal warnings so
# Cloud logs remain useful for scan diagnostics.
warnings.filterwarnings(
    "ignore",
    message=r"Data Validation extension is not supported and will be removed",
    category=UserWarning,
    module=r"openpyxl\..*",
)
warnings.filterwarnings(
    "ignore",
    message=r"Unknown extension is not supported and will be removed",
    category=UserWarning,
    module=r"openpyxl\..*",
)
warnings.filterwarnings(
    "ignore",
    message=r"Workbook contains no default style, apply openpyxl's default",
    category=UserWarning,
    module=r"openpyxl\..*",
)


TEXT_EXTENSIONS = {
    ".pdf",
    ".docx",
    ".pptx",
    ".xlsx",
    ".csv",
    ".txt",
    ".md",
    ".html",
    ".htm",
    ".xml",
    ".json",
    ".log",
}

IMAGE_EXTENSIONS = {
    ".png",
    ".jpg",
    ".jpeg",
    ".webp",
    ".bmp",
    ".tif",
    ".tiff",
}

SUPPORTED_EXTENSIONS = TEXT_EXTENSIONS | IMAGE_EXTENSIONS
UPLOAD_EXTENSIONS = SUPPORTED_EXTENSIONS | {".zip"}

DEFAULT_MIN_FILES = 10
DEFAULT_MIN_SOURCE_TOKENS = 256_000
DEFAULT_MAX_CONTEXT_TOKENS = 1_000_000

# Generic tiled-vision proxy. These constants intentionally live in one place
# so they can be changed if the project confirms a model-specific image rule.
VISION_BASE_TOKENS = 85
VISION_TILE_TOKENS = 170
VISION_TILE_SIZE = 512
VISION_MAX_DIMENSION = 2048
VISION_TARGET_SHORT_SIDE = 768
MIN_RASTER_DIMENSION = 64

MAX_ZIP_MEMBERS = 2_000
# The engine keeps the original 3 GiB total ZIP allowance so environment coverage is not
# reduced. Extraction is now streamed, so the safety limit is primarily a zip-bomb/disk guard.
MAX_ZIP_UNCOMPRESSED_BYTES = 3 * 1024 * 1024 * 1024
MAX_ZIP_MEMBER_BYTES = 2 * 1024 * 1024 * 1024
MAX_ZIP_COMPRESSION_RATIO = 500.0
STREAM_COPY_CHUNK_BYTES = 1024 * 1024
TOKEN_CHUNK_CHARACTERS = 250_000
LARGE_JSON_STREAM_THRESHOLD_BYTES = 32 * 1024 * 1024


@dataclass
class FileEstimate:
    file: str
    extension: str
    size_mb: float
    characters: int
    words: int
    text_tokens: int
    image_tokens: int
    estimated_tokens: int
    image_count: int
    image_estimate_confidence: str
    maximum_risk: str
    extraction_notes: str


def clean_text(text: str) -> str:
    """Normalize extracted text before estimating tokens."""
    if not text:
        return ""
    text = text.replace("\x00", " ")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


@lru_cache(maxsize=1)
def _get_token_encoder():
    """Return a cached tiktoken encoder, or None when tiktoken is unavailable."""
    try:
        import tiktoken  # type: ignore

        return tiktoken.get_encoding("cl100k_base")
    except Exception:
        return None


def _iter_bounded_text(text: str, limit: int = TOKEN_CHUNK_CHARACTERS):
    """Yield bounded pieces so a single tokenizer call cannot create a huge token list."""
    if not text:
        return
    for start in range(0, len(text), limit):
        yield text[start : start + limit]


def _encode_ordinary_token_count(encoder, text: str) -> int:
    """Count source-document text as ordinary text, even when it resembles model special tokens.

    Source files can legitimately contain literal strings such as ``<|endoftext|>``. Those
    strings are document content, not tokenizer control instructions, so the estimator must not
    reject them as disallowed special tokens. ``encode_ordinary`` gives the desired behavior;
    the fallback keeps compatibility with encoder implementations that expose only ``encode``.
    """
    if not text:
        return 0
    encode_ordinary = getattr(encoder, "encode_ordinary", None)
    if callable(encode_ordinary):
        return len(encode_ordinary(text))
    return len(encoder.encode(text, disallowed_special=()))


@dataclass
class TextMetrics:
    characters: int
    words: int
    text_tokens: int
    extraction_notes: str


class _TextMetricsAccumulator:
    """Incrementally count cleaned text without retaining the full extracted document."""

    def __init__(self) -> None:
        self.characters = 0
        self.words = 0
        self.text_tokens = 0
        self._has_text = False
        self._encoder = _get_token_encoder()
        self._token_buffer: list[str] = []
        self._token_buffer_characters = 0

    def _flush_token_buffer(self) -> None:
        if self._encoder is None or not self._token_buffer:
            return
        text = "".join(self._token_buffer)
        self.text_tokens += _encode_ordinary_token_count(self._encoder, text)
        self._token_buffer.clear()
        self._token_buffer_characters = 0

    def _add_for_tokenization(self, piece: str) -> None:
        if self._encoder is None:
            return
        # Normal logical chunks (rows/pages/paragraphs) are batched together for speed and
        # to preserve BPE behavior across most boundaries. A pathological giant row/string is
        # split by itself so one encode call still stays bounded.
        if len(piece) > TOKEN_CHUNK_CHARACTERS:
            self._flush_token_buffer()
            for bounded in _iter_bounded_text(piece):
                self.text_tokens += _encode_ordinary_token_count(self._encoder, bounded)
            return

        if (
            self._token_buffer
            and self._token_buffer_characters + len(piece) > TOKEN_CHUNK_CHARACTERS
        ):
            self._flush_token_buffer()
        self._token_buffer.append(piece)
        self._token_buffer_characters += len(piece)

    def add(self, text: str) -> None:
        cleaned = clean_text(text)
        if not cleaned:
            return

        piece = ("\n" if self._has_text else "") + cleaned
        self._has_text = True
        self.characters += len(piece)
        self.words += len(re.findall(r"\S+", piece))
        self._add_for_tokenization(piece)

    def finish(self, notes: str) -> TextMetrics:
        self._flush_token_buffer()
        tokens = self.text_tokens
        if self._encoder is None and self.characters:
            tokens = int(max(self.characters / 4.0, self.words * 1.35))
        return TextMetrics(
            characters=self.characters,
            words=self.words,
            text_tokens=tokens,
            extraction_notes=notes,
        )


class _VisibleHTMLAccumulator(HTMLParser):
    """Streaming HTML text collector that ignores script/style/noscript contents."""

    def __init__(self, accumulator: _TextMetricsAccumulator) -> None:
        super().__init__(convert_charrefs=True)
        self.accumulator = accumulator
        self._ignored_depth = 0

    def handle_starttag(self, tag: str, attrs) -> None:
        if tag.lower() in {"script", "style", "noscript"}:
            self._ignored_depth += 1

    def handle_endtag(self, tag: str) -> None:
        if tag.lower() in {"script", "style", "noscript"} and self._ignored_depth:
            self._ignored_depth -= 1

    def handle_data(self, data: str) -> None:
        if not self._ignored_depth and data.strip():
            self.accumulator.add(data)


def estimate_text_tokens(text: str) -> int:
    """Estimate text tokens using tiktoken when available, else a heuristic."""
    text = text or ""
    if not text:
        return 0

    encoder = _get_token_encoder()
    if encoder is not None:
        # Bound each encode call to avoid allocating one enormous token list.
        return sum(_encode_ordinary_token_count(encoder, piece) for piece in _iter_bounded_text(text))

    characters = len(text)
    words = len(re.findall(r"\S+", text))
    return int(max(characters / 4.0, words * 1.35))


# Backward-compatible alias for code that imported estimate_tokens.
estimate_tokens = estimate_text_tokens


def _normalized_visual_dimensions(width: int, height: int) -> tuple[int, int]:
    """Normalize raster dimensions for the generic tiled-vision proxy."""
    width = max(int(width), 1)
    height = max(int(height), 1)

    max_side = max(width, height)
    if max_side > VISION_MAX_DIMENSION:
        scale = VISION_MAX_DIMENSION / max_side
        width = max(1, int(round(width * scale)))
        height = max(1, int(round(height * scale)))

    short_side = min(width, height)
    if short_side > VISION_TARGET_SHORT_SIDE:
        scale = VISION_TARGET_SHORT_SIDE / short_side
        width = max(1, int(round(width * scale)))
        height = max(1, int(round(height * scale)))

    return width, height


def estimate_image_tokens_from_dimensions(width: int, height: int) -> int:
    """
    Estimate visual tokens from raster dimensions.

    This is deliberately a lower-confidence proxy, not a claim about the
    target model's exact vision tokenizer. It approximates tiled high-detail
    processing after conservative image normalization.
    """
    if width < MIN_RASTER_DIMENSION or height < MIN_RASTER_DIMENSION:
        return 0

    width, height = _normalized_visual_dimensions(width, height)
    tiles = math.ceil(width / VISION_TILE_SIZE) * math.ceil(height / VISION_TILE_SIZE)
    return VISION_BASE_TOKENS + tiles * VISION_TILE_TOKENS


def _image_dimensions_from_fileobj(fileobj) -> Optional[tuple[int, int]]:
    try:
        from PIL import Image

        with Image.open(fileobj) as image:
            return int(image.width), int(image.height)
    except Exception:
        return None


def _office_media_estimate(path: Path) -> tuple[int, int, str]:
    """Estimate raster images in Office packages without reading full images into RAM."""
    tokens = 0
    count = 0
    try:
        with zipfile.ZipFile(path) as archive:
            for info in archive.infolist():
                member = info.filename.lower()
                if "/media/" not in member:
                    continue
                if Path(member).suffix.lower() not in IMAGE_EXTENSIONS:
                    continue
                try:
                    with archive.open(info) as source:
                        dimensions = _image_dimensions_from_fileobj(source)
                except Exception:
                    dimensions = None
                if not dimensions:
                    continue
                image_tokens = estimate_image_tokens_from_dimensions(*dimensions)
                if image_tokens:
                    tokens += image_tokens
                    count += 1
    except Exception as exc:
        return 0, 0, f"Embedded-image scan failed: {exc}"

    if count:
        return tokens, count, f"Embedded raster images estimated: {count}"
    return 0, 0, "No qualifying embedded raster images detected"


def _pdf_media_estimate(path: Path) -> tuple[int, int, str]:
    """Estimate embedded PDF raster images from metadata, deduplicating image xrefs."""
    try:
        import pymupdf as fitz
    except Exception as exc:
        return 0, 0, f"PDF image scan skipped: PyMuPDF unavailable ({exc})"

    tokens = 0
    count = 0
    seen_xrefs: set[int] = set()

    try:
        with fitz.open(path) as document:
            for page in document:
                for image in page.get_images(full=True):
                    xref = int(image[0])
                    if xref in seen_xrefs:
                        continue
                    seen_xrefs.add(xref)
                    width = int(image[2] or 0)
                    height = int(image[3] or 0)
                    image_tokens = estimate_image_tokens_from_dimensions(width, height)
                    if image_tokens:
                        tokens += image_tokens
                        count += 1
    except Exception as exc:
        return 0, 0, f"PDF image scan failed: {exc}"

    if count:
        return tokens, count, f"Embedded PDF raster images estimated: {count}"
    return 0, 0, "No qualifying embedded PDF raster images detected"


def estimate_visual_tokens(path: Path) -> tuple[int, int, str, str]:
    """Return image tokens, image count, note, and confidence label."""
    extension = path.suffix.lower()

    if extension in IMAGE_EXTENSIONS:
        try:
            from PIL import Image

            with Image.open(path) as image:
                tokens = estimate_image_tokens_from_dimensions(image.width, image.height)
                count = 1 if tokens else 0
                return tokens, count, f"Standalone raster: {image.width}×{image.height}", "LOWER"
        except Exception as exc:
            return 0, 0, f"Standalone image estimate failed: {exc}", "LOWER"

    if extension == ".pdf":
        tokens, count, note = _pdf_media_estimate(path)
        return tokens, count, note, "LOWER"

    if extension in {".docx", ".pptx", ".xlsx"}:
        tokens, count, note = _office_media_estimate(path)
        return tokens, count, note, "LOWER"

    return 0, 0, "Visual estimate not applicable for this format", "N/A"


def extract_pdf(path: Path, max_pages: Optional[int] = None) -> tuple[str, str]:
    """Extract text from a PDF."""
    try:
        import pymupdf as fitz
    except Exception as exc:
        return "", f"PDF skipped: PyMuPDF not installed ({exc})"

    notes: list[str] = []
    chunks: list[str] = []

    try:
        with fitz.open(path) as document:
            page_count = len(document)
            pages_to_read = page_count if max_pages is None else min(page_count, max_pages)

            for page_index in range(pages_to_read):
                chunks.append(document.load_page(page_index).get_text("text"))

            notes.append(f"PDF pages read: {pages_to_read}/{page_count}")
            if max_pages is not None and page_count > max_pages:
                notes.append("PDF truncated by max-pages option")
    except Exception as exc:
        return "", f"PDF extraction failed: {exc}"

    return clean_text("\n".join(chunks)), "; ".join(notes)


def extract_docx(path: Path) -> tuple[str, str]:
    """Extract paragraphs and tables from a DOCX file."""
    try:
        import docx
    except Exception as exc:
        return "", f"DOCX skipped: python-docx not installed ({exc})"

    chunks: list[str] = []
    try:
        document = docx.Document(str(path))
        for paragraph in document.paragraphs:
            if paragraph.text:
                chunks.append(paragraph.text)
        for table in document.tables:
            for row in table.rows:
                chunks.append("\t".join(cell.text for cell in row.cells))

        notes = f"DOCX paragraphs: {len(document.paragraphs)}, tables: {len(document.tables)}"
        return clean_text("\n".join(chunks)), notes
    except Exception as exc:
        return "", f"DOCX extraction failed: {exc}"


def extract_pptx_text(path: Path) -> tuple[str, str]:
    """Extract slide text and speaker notes from a PPTX file."""
    try:
        presentation = Presentation(str(path))
        parts: list[str] = []
        for slide_index, slide in enumerate(presentation.slides, start=1):
            parts.append(f"\n--- Slide {slide_index} ---\n")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    parts.append(shape.text)
            if slide.has_notes_slide:
                try:
                    notes_text_frame = slide.notes_slide.notes_text_frame
                    if notes_text_frame and notes_text_frame.text:
                        parts.append(f"\n[Speaker notes]\n{notes_text_frame.text}")
                except Exception:
                    pass
        return clean_text("\n".join(parts)), f"PPTX slides read: {len(presentation.slides)}"
    except Exception as exc:
        return "", f"PPTX extraction failed: {exc}"


def extract_xlsx(path: Path, max_cells: Optional[int] = None) -> tuple[str, str]:
    """Extract displayed cell values and formulas from an XLSX file."""
    try:
        from openpyxl import load_workbook
    except Exception as exc:
        return "", f"XLSX skipped: openpyxl not installed ({exc})"

    chunks: list[str] = []
    cell_count = 0
    try:
        workbook = load_workbook(filename=path, read_only=True, data_only=False)
        sheet_count = len(workbook.sheetnames)

        for worksheet in workbook.worksheets:
            chunks.append(f"\n[SHEET: {worksheet.title}]\n")
            for row in worksheet.iter_rows():
                values: list[str] = []
                for cell in row:
                    if cell.value is not None:
                        values.append(str(cell.value))
                        cell_count += 1
                        if max_cells is not None and cell_count >= max_cells:
                            if values:
                                chunks.append("\t".join(values))
                            workbook.close()
                            notes = (
                                f"XLSX sheets: {sheet_count}; cells read: {cell_count}; "
                                "truncated by max-cells option"
                            )
                            return clean_text("\n".join(chunks)), notes
                if values:
                    chunks.append("\t".join(values))

        workbook.close()
        return clean_text("\n".join(chunks)), (
            f"XLSX sheets: {sheet_count}; non-empty cells read: {cell_count}"
        )
    except Exception as exc:
        return "", f"XLSX extraction failed: {exc}"


def extract_csv(path: Path, max_rows: Optional[int] = None) -> tuple[str, str]:
    """Extract rows from a CSV or delimited-text file."""
    chunks: list[str] = []
    row_count = 0
    try:
        with path.open("r", encoding="utf-8-sig", errors="replace", newline="") as file:
            sample = file.read(4096)
            file.seek(0)
            try:
                dialect = csv.Sniffer().sniff(sample)
            except Exception:
                dialect = csv.excel

            reader = csv.reader(file, dialect)
            for row in reader:
                row_count += 1
                chunks.append("\t".join(row))
                if max_rows is not None and row_count >= max_rows:
                    return clean_text("\n".join(chunks)), (
                        f"CSV rows read: {row_count}; truncated by max-rows option"
                    )

        return clean_text("\n".join(chunks)), f"CSV rows read: {row_count}"
    except Exception as exc:
        return "", f"CSV extraction failed: {exc}"


def extract_html_text(path: Path) -> tuple[str, str]:
    """Extract visible text from an HTML file."""
    try:
        raw = path.read_text(encoding="utf-8", errors="ignore")
        soup = BeautifulSoup(raw, "html.parser")
        for tag in soup(["script", "style", "noscript"]):
            tag.decompose()
        return clean_text(soup.get_text(separator="\n", strip=True)), (
            "HTML visible text extracted; scripts/styles removed"
        )
    except Exception as exc:
        return "", f"HTML extraction failed: {exc}"


def extract_xml_text(path: Path) -> tuple[str, str]:
    """Extract tags and text values from XML."""
    try:
        raw = path.read_text(encoding="utf-8", errors="ignore")
    except Exception as exc:
        return "", f"XML read failed: {exc}"

    try:
        root = ET.fromstring(raw)
    except ET.ParseError:
        return clean_text(raw), "XML parse failed; raw text used"

    parts: list[str] = []

    def walk(node, depth: int = 0) -> None:
        if node.tag:
            parts.append(f"{'  ' * depth}<{node.tag}>")
        if node.text and node.text.strip():
            parts.append(node.text.strip())
        for child in node:
            walk(child, depth + 1)
        if node.tail and node.tail.strip():
            parts.append(node.tail.strip())

    walk(root)
    return clean_text("\n".join(parts)), "XML text extracted from element tree"


def extract_json_text(path: Path) -> tuple[str, str]:
    """Read JSON as normalized textual content without changing values."""
    try:
        raw = path.read_text(encoding="utf-8", errors="replace")
        try:
            parsed = json.loads(raw)
            text = json.dumps(parsed, ensure_ascii=False, separators=(",", ":"))
            return clean_text(text), "JSON parsed and normalized"
        except Exception:
            return clean_text(raw), "JSON parse failed; raw text used"
    except Exception as exc:
        return "", f"JSON extraction failed: {exc}"


def extract_text(path: Path) -> tuple[str, str]:
    """Route a file to the appropriate text extraction function."""
    extension = path.suffix.lower()

    if extension == ".pdf":
        return extract_pdf(path)
    if extension == ".docx":
        return extract_docx(path)
    if extension == ".pptx":
        return extract_pptx_text(path)
    if extension == ".xlsx":
        return extract_xlsx(path)
    if extension == ".csv":
        return extract_csv(path)
    if extension in {".txt", ".md", ".log"}:
        try:
            return clean_text(path.read_text(encoding="utf-8", errors="replace")), (
                "Plain text/Markdown/log read"
            )
        except Exception as exc:
            return "", f"Text extraction failed: {exc}"
    if extension in {".html", ".htm"}:
        return extract_html_text(path)
    if extension == ".xml":
        return extract_xml_text(path)
    if extension == ".json":
        return extract_json_text(path)
    if extension in IMAGE_EXTENSIONS:
        return "", "Standalone image: no OCR text added"

    return "", f"Unsupported extension: {extension}"


def estimate_text_metrics(path: Path) -> TextMetrics:
    """Extract text incrementally and return counts without retaining full document text."""
    extension = path.suffix.lower()
    accumulator = _TextMetricsAccumulator()

    if extension in IMAGE_EXTENSIONS:
        return accumulator.finish("Standalone image: no OCR text added")

    if extension == ".pdf":
        try:
            import pymupdf as fitz
        except Exception as exc:
            return accumulator.finish(f"PDF skipped: PyMuPDF not installed ({exc})")
        try:
            with fitz.open(path) as document:
                page_count = len(document)
                for page in document:
                    accumulator.add(page.get_text("text"))
            return accumulator.finish(f"PDF pages read: {page_count}/{page_count}; streamed metrics")
        except Exception as exc:
            return accumulator.finish(f"PDF extraction failed: {exc}")

    if extension == ".docx":
        try:
            import docx
        except Exception as exc:
            return accumulator.finish(f"DOCX skipped: python-docx not installed ({exc})")
        try:
            document = docx.Document(str(path))
            for paragraph in document.paragraphs:
                if paragraph.text:
                    accumulator.add(paragraph.text)
            for table in document.tables:
                for row in table.rows:
                    accumulator.add("\t".join(cell.text for cell in row.cells))
            return accumulator.finish(
                f"DOCX paragraphs: {len(document.paragraphs)}, tables: {len(document.tables)}; streamed metrics"
            )
        except Exception as exc:
            return accumulator.finish(f"DOCX extraction failed: {exc}")

    if extension == ".pptx":
        try:
            presentation = Presentation(str(path))
            slide_count = len(presentation.slides)
            for slide_index, slide in enumerate(presentation.slides, start=1):
                accumulator.add(f"--- Slide {slide_index} ---")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        accumulator.add(shape.text)
                if slide.has_notes_slide:
                    try:
                        notes_text_frame = slide.notes_slide.notes_text_frame
                        if notes_text_frame and notes_text_frame.text:
                            accumulator.add(f"[Speaker notes]\n{notes_text_frame.text}")
                    except Exception:
                        pass
            return accumulator.finish(f"PPTX slides read: {slide_count}; streamed metrics")
        except Exception as exc:
            return accumulator.finish(f"PPTX extraction failed: {exc}")

    if extension == ".xlsx":
        try:
            from openpyxl import load_workbook
        except Exception as exc:
            return accumulator.finish(f"XLSX skipped: openpyxl not installed ({exc})")
        cell_count = 0
        try:
            workbook = load_workbook(filename=path, read_only=True, data_only=False)
            sheet_count = len(workbook.sheetnames)
            try:
                for worksheet in workbook.worksheets:
                    accumulator.add(f"[SHEET: {worksheet.title}]")
                    for row in worksheet.iter_rows():
                        values = []
                        for cell in row:
                            if cell.value is not None:
                                values.append(str(cell.value))
                                cell_count += 1
                        if values:
                            accumulator.add("\t".join(values))
            finally:
                workbook.close()
            return accumulator.finish(
                f"XLSX sheets: {sheet_count}; non-empty cells read: {cell_count}; streamed metrics"
            )
        except Exception as exc:
            return accumulator.finish(f"XLSX extraction failed: {exc}")

    if extension == ".csv":
        row_count = 0
        try:
            with path.open("r", encoding="utf-8-sig", errors="replace", newline="") as file:
                sample = file.read(4096)
                file.seek(0)
                try:
                    dialect = csv.Sniffer().sniff(sample)
                except Exception:
                    dialect = csv.excel
                reader = csv.reader(file, dialect)
                for row in reader:
                    row_count += 1
                    accumulator.add("\t".join(row))
            return accumulator.finish(f"CSV rows read: {row_count}; streamed metrics")
        except Exception as exc:
            return accumulator.finish(f"CSV extraction failed: {exc}")

    if extension in {".txt", ".md", ".log"}:
        try:
            with path.open("r", encoding="utf-8", errors="replace") as file:
                for line in file:
                    accumulator.add(line)
            return accumulator.finish("Plain text/Markdown/log read incrementally")
        except Exception as exc:
            return accumulator.finish(f"Text extraction failed: {exc}")

    if extension in {".html", ".htm"}:
        try:
            parser = _VisibleHTMLAccumulator(accumulator)
            with path.open("r", encoding="utf-8", errors="ignore") as file:
                while True:
                    chunk = file.read(1024 * 1024)
                    if not chunk:
                        break
                    parser.feed(chunk)
            parser.close()
            return accumulator.finish("HTML visible text extracted incrementally; scripts/styles removed")
        except Exception as exc:
            return accumulator.finish(f"HTML extraction failed: {exc}")

    if extension == ".xml":
        try:
            # iterparse keeps memory bounded by clearing elements once consumed.
            for _, elem in ET.iterparse(path, events=("end",)):
                if elem.tag:
                    accumulator.add(f"<{elem.tag}>")
                if elem.text and elem.text.strip():
                    accumulator.add(elem.text.strip())
                if elem.tail and elem.tail.strip():
                    accumulator.add(elem.tail.strip())
                elem.clear()
            return accumulator.finish("XML text extracted incrementally from element tree")
        except ET.ParseError:
            try:
                with path.open("r", encoding="utf-8", errors="ignore") as file:
                    for line in file:
                        accumulator.add(line)
                return accumulator.finish("XML parse failed; raw text read incrementally")
            except Exception as exc:
                return accumulator.finish(f"XML extraction failed: {exc}")
        except Exception as exc:
            return accumulator.finish(f"XML extraction failed: {exc}")

    if extension == ".json":
        try:
            file_size = path.stat().st_size
        except Exception:
            file_size = 0
        if file_size and file_size <= LARGE_JSON_STREAM_THRESHOLD_BYTES:
            try:
                raw = path.read_text(encoding="utf-8", errors="replace")
                try:
                    parsed = json.loads(raw)
                    normalized = json.dumps(parsed, ensure_ascii=False, separators=(",", ":"))
                    accumulator.add(normalized)
                    return accumulator.finish("JSON parsed and normalized")
                except Exception:
                    accumulator.add(raw)
                    return accumulator.finish("JSON parse failed; raw text used")
            except Exception as exc:
                return accumulator.finish(f"JSON extraction failed: {exc}")
        try:
            with path.open("r", encoding="utf-8", errors="replace") as file:
                while True:
                    chunk = file.read(1024 * 1024)
                    if not chunk:
                        break
                    accumulator.add(chunk)
            return accumulator.finish(
                "Large JSON read incrementally; normalization skipped to keep memory bounded"
            )
        except Exception as exc:
            return accumulator.finish(f"JSON extraction failed: {exc}")

    return accumulator.finish(f"Unsupported extension: {extension}")


def maximum_risk_band(tokens: int, maximum_tokens: int = DEFAULT_MAX_CONTEXT_TOKENS) -> str:
    """Assign a maximum-context risk label."""
    if maximum_tokens <= 0:
        raise ValueError("The maximum context limit must be greater than zero.")

    ratio = tokens / maximum_tokens
    if ratio >= 1.0:
        return "OVER_MAXIMUM"
    if ratio >= 0.95:
        return "CRITICAL"
    if ratio >= 0.85:
        return "NEAR_MAXIMUM"
    if ratio >= 0.70:
        return "HIGH_CONTEXT_LOAD"
    return "COMFORTABLY_WITHIN_MAXIMUM"


def estimate_file(
    path: Path,
    maximum_tokens: int = DEFAULT_MAX_CONTEXT_TOKENS,
    display_name: Optional[str] = None,
) -> FileEstimate:
    """Extract and estimate text + visual token usage for one file with bounded memory."""
    metrics = estimate_text_metrics(path)
    visual_tokens, image_count, visual_note, confidence = estimate_visual_tokens(path)

    total_tokens = metrics.text_tokens + visual_tokens

    try:
        size_mb = path.stat().st_size / (1024 * 1024)
    except Exception:
        size_mb = 0.0

    notes = (
        f"{metrics.extraction_notes}; {visual_note}"
        if visual_note
        else metrics.extraction_notes
    )

    return FileEstimate(
        file=display_name or str(path),
        extension=path.suffix.lower(),
        size_mb=round(size_mb, 2),
        characters=metrics.characters,
        words=metrics.words,
        text_tokens=metrics.text_tokens,
        image_tokens=visual_tokens,
        estimated_tokens=total_tokens,
        image_count=image_count,
        image_estimate_confidence=confidence,
        maximum_risk=maximum_risk_band(total_tokens, maximum_tokens),
        extraction_notes=notes,
    )

def estimate_prompt(prompt: str) -> FileEstimate:
    """Estimate prompt text separately; prompt tokens never satisfy the 256k source minimum."""
    text = prompt or ""
    characters = len(text)
    words = len(re.findall(r"\S+", text))
    tokens = estimate_text_tokens(text)

    return FileEstimate(
        file="[PROMPT_TEXT]",
        extension="prompt",
        size_mb=0.0,
        characters=characters,
        words=words,
        text_tokens=tokens,
        image_tokens=0,
        estimated_tokens=tokens,
        image_count=0,
        image_estimate_confidence="N/A",
        maximum_risk="N/A",
        extraction_notes="Prompt text supplied manually; excluded from source minimum",
    )


def summarize_project(
    required_estimates: Sequence[FileEstimate],
    all_estimates: Optional[Sequence[FileEstimate]] = None,
    prompt_estimate: Optional[FileEstimate] = None,
    min_files: int = DEFAULT_MIN_FILES,
    min_source_tokens: int = DEFAULT_MIN_SOURCE_TOKENS,
    max_context_tokens: int = DEFAULT_MAX_CONTEXT_TOKENS,
) -> dict:
    """Summarize minimum compliance and maximum-context safety separately."""
    if min_files <= 0:
        raise ValueError("Minimum file count must be greater than zero.")
    if min_source_tokens <= 0:
        raise ValueError("Minimum source tokens must be greater than zero.")
    if max_context_tokens <= 0:
        raise ValueError("Maximum context tokens must be greater than zero.")
    if max_context_tokens <= min_source_tokens:
        raise ValueError("Maximum context must be greater than the minimum source-token threshold.")

    required_estimates = list(required_estimates)
    all_estimates = list(all_estimates if all_estimates is not None else required_estimates)

    required_file_count = len(required_estimates)
    uploaded_file_count = len(all_estimates)
    required_text_tokens = sum(item.text_tokens for item in required_estimates)
    required_image_tokens = sum(item.image_tokens for item in required_estimates)
    required_source_tokens = required_text_tokens + required_image_tokens

    all_text_tokens = sum(item.text_tokens for item in all_estimates)
    all_image_tokens = sum(item.image_tokens for item in all_estimates)
    all_source_tokens = all_text_tokens + all_image_tokens
    prompt_tokens = prompt_estimate.estimated_tokens if prompt_estimate else 0

    required_turn_tokens = required_source_tokens + prompt_tokens
    full_uploaded_turn_tokens = all_source_tokens + prompt_tokens

    files_met = required_file_count >= min_files
    source_met = required_source_tokens >= min_source_tokens
    minimum_met = files_met and source_met

    minimum_ratio = required_source_tokens / min_source_tokens
    if not files_met and not source_met:
        minimum_status = "BELOW_MINIMUM_FILES_AND_CONTEXT"
    elif not files_met:
        minimum_status = "BELOW_MINIMUM_FILES"
    elif not source_met:
        minimum_status = "BELOW_MINIMUM_CONTEXT"
    elif minimum_ratio < 1.10:
        minimum_status = "MEETS_MINIMUM_NARROW_MARGIN"
    else:
        minimum_status = "COMFORTABLY_ABOVE_MINIMUM"

    max_status = maximum_risk_band(full_uploaded_turn_tokens, max_context_tokens)
    if not minimum_met:
        overall_status = "BELOW_PROJECT_MINIMUM"
    elif max_status == "OVER_MAXIMUM":
        overall_status = "OVER_SELECTED_MAXIMUM"
    elif max_status in {"CRITICAL", "NEAR_MAXIMUM"}:
        overall_status = "MEETS_MINIMUM_BUT_NEAR_MAXIMUM"
    else:
        overall_status = "MEETS_PROJECT_REQUIREMENTS"

    required_file_shortfall = max(0, min_files - required_file_count)
    source_token_shortfall = max(0, min_source_tokens - required_source_tokens)
    remaining_to_maximum = max(0, max_context_tokens - full_uploaded_turn_tokens)

    utilization_pct = (
        required_file_count / uploaded_file_count * 100 if uploaded_file_count else 0.0
    )
    environment_utilization_advisory_met = (
        required_file_count >= 20 or utilization_pct >= 50.0
    ) if uploaded_file_count else False

    largest_required = sorted(
        required_estimates, key=lambda item: item.estimated_tokens, reverse=True
    )[:10]
    largest_uploaded = sorted(
        all_estimates, key=lambda item: item.estimated_tokens, reverse=True
    )[:10]

    return {
        "min_files": min_files,
        "min_source_tokens": min_source_tokens,
        "max_context_tokens": max_context_tokens,
        "required_file_count": required_file_count,
        "uploaded_file_count": uploaded_file_count,
        "files_met": files_met,
        "required_file_shortfall": required_file_shortfall,
        "required_text_tokens": required_text_tokens,
        "required_image_tokens": required_image_tokens,
        "required_source_tokens": required_source_tokens,
        "source_met": source_met,
        "source_token_shortfall": source_token_shortfall,
        "minimum_status": minimum_status,
        "minimum_met": minimum_met,
        "all_text_tokens": all_text_tokens,
        "all_image_tokens": all_image_tokens,
        "all_source_tokens": all_source_tokens,
        "prompt_tokens": prompt_tokens,
        "required_turn_tokens": required_turn_tokens,
        "full_uploaded_turn_tokens": full_uploaded_turn_tokens,
        "percent_of_maximum": round(full_uploaded_turn_tokens / max_context_tokens * 100, 1),
        "remaining_to_maximum": remaining_to_maximum,
        "maximum_status": max_status,
        "overall_status": overall_status,
        "required_file_utilization_pct": round(utilization_pct, 1),
        "environment_utilization_advisory_met": environment_utilization_advisory_met,
        "largest_required_contributors": [asdict(item) for item in largest_required],
        "largest_uploaded_contributors": [asdict(item) for item in largest_uploaded],
    }


# Backward-compatible summary for simple callers: all files are required.
def summarize(
    estimates: List[FileEstimate],
    context_limit: int = DEFAULT_MAX_CONTEXT_TOKENS,
) -> dict:
    return summarize_project(
        required_estimates=estimates,
        all_estimates=estimates,
        max_context_tokens=context_limit,
    )


def iter_files(paths: Iterable[Path]) -> Iterable[Path]:
    """Yield supported non-ZIP files from individual paths or folders."""
    for path in paths:
        if path.is_dir():
            for candidate in sorted(path.rglob("*")):
                if candidate.is_file() and candidate.suffix.lower() in SUPPORTED_EXTENSIONS:
                    yield candidate
        elif path.is_file() and path.suffix.lower() in SUPPORTED_EXTENSIONS:
            yield path


def inspect_zip(zip_path: Path) -> dict:
    """Return ZIP safety/diagnostic metadata without extracting file contents."""
    with zipfile.ZipFile(zip_path) as archive:
        infos = archive.infolist()
        member_count = len(infos)
        if member_count > MAX_ZIP_MEMBERS:
            raise ValueError(f"ZIP has {member_count} members; limit is {MAX_ZIP_MEMBERS}.")

        total_uncompressed = sum(max(0, int(info.file_size)) for info in infos)
        total_compressed = sum(max(0, int(info.compress_size)) for info in infos)
        if total_uncompressed > MAX_ZIP_UNCOMPRESSED_BYTES:
            raise ValueError(
                f"ZIP uncompressed size is {total_uncompressed / (1024**3):.2f} GiB; "
                f"limit is {MAX_ZIP_UNCOMPRESSED_BYTES / (1024**3):.2f} GiB."
            )

        largest_member = 0
        max_ratio = 0.0
        supported_count = 0
        for info in infos:
            if info.is_dir():
                continue
            largest_member = max(largest_member, int(info.file_size))
            if int(info.file_size) > MAX_ZIP_MEMBER_BYTES:
                raise ValueError(
                    f"ZIP member {info.filename!r} is {info.file_size / (1024**3):.2f} GiB; "
                    f"per-member limit is {MAX_ZIP_MEMBER_BYTES / (1024**3):.2f} GiB."
                )
            if info.file_size:
                ratio = info.file_size / max(1, info.compress_size)
                max_ratio = max(max_ratio, ratio)
                if ratio > MAX_ZIP_COMPRESSION_RATIO and info.file_size > 50 * 1024 * 1024:
                    raise ValueError(
                        f"ZIP member {info.filename!r} has an unsafe compression ratio "
                        f"({ratio:.0f}:1)."
                    )
            if Path(PurePosixPath(info.filename).name).suffix.lower() in SUPPORTED_EXTENSIONS:
                supported_count += 1

        return {
            "members": member_count,
            "supported_members": supported_count,
            "compressed_bytes": total_compressed,
            "uncompressed_bytes": total_uncompressed,
            "largest_member_bytes": largest_member,
            "maximum_compression_ratio": round(max_ratio, 1),
        }


def extract_supported_from_zip(zip_path: Path, destination: Path) -> list[tuple[Path, str]]:
    """
    Safely extract supported files from a ZIP using bounded-memory streaming copies.

    Returns (local_path, archive_relative_name) pairs. Unsupported members,
    directories, path traversal entries, and macOS metadata are ignored.
    """
    extracted: list[tuple[Path, str]] = []
    destination.mkdir(parents=True, exist_ok=True)
    inspect_zip(zip_path)  # validates member count, sizes, and compression ratios first

    with zipfile.ZipFile(zip_path) as archive:
        for info in archive.infolist():
            if info.is_dir():
                continue

            pure = PurePosixPath(info.filename)
            if pure.is_absolute() or ".." in pure.parts:
                continue
            if "__MACOSX" in pure.parts or pure.name.startswith("._"):
                continue

            extension = Path(pure.name).suffix.lower()
            if extension not in SUPPORTED_EXTENSIONS:
                continue

            target = destination.joinpath(*pure.parts)
            target.parent.mkdir(parents=True, exist_ok=True)
            with archive.open(info) as source, target.open("wb") as output:
                shutil.copyfileobj(source, output, length=STREAM_COPY_CHUNK_BYTES)

            extracted.append((target, str(pure)))

    return extracted

def print_table(
    estimates: Sequence[FileEstimate],
    summary: dict,
) -> None:
    """Print a CLI report."""
    print("\nProfessional Tasks context range estimate")
    print("=" * 112)
    print(
        f"{'Total':>11}  {'Text':>11}  {'Image':>10}  {'Words':>10}  "
        f"{'MB':>8}  {'Max risk':<29}  File"
    )
    print("-" * 112)

    for estimate in sorted(estimates, key=lambda item: item.estimated_tokens, reverse=True):
        print(
            f"{estimate.estimated_tokens:>11,}  "
            f"{estimate.text_tokens:>11,}  "
            f"{estimate.image_tokens:>10,}  "
            f"{estimate.words:>10,}  "
            f"{estimate.size_mb:>8.2f}  "
            f"{estimate.maximum_risk:<29}  "
            f"{estimate.file}"
        )

    print("-" * 112)
    print(
        f"REQUIRED SOURCE: {summary['required_source_tokens']:,} / "
        f"{summary['min_source_tokens']:,} tokens | "
        f"FILES: {summary['required_file_count']} / {summary['min_files']}"
    )
    print(
        f"FULL UPLOADED TURN: {summary['full_uploaded_turn_tokens']:,} / "
        f"{summary['max_context_tokens']:,} tokens "
        f"({summary['percent_of_maximum']}%)"
    )
    print(f"OVERALL STATUS: {summary['overall_status']}")
    print(
        "\nImage-token estimates are a lower-confidence raster-dimension proxy. "
        "They are reported separately from text tokens and are not OCR-derived."
    )


def _collect_cli_files(paths: Sequence[Path], tmp_root: Path) -> list[tuple[Path, str]]:
    collected: list[tuple[Path, str]] = []
    for path in paths:
        if path.is_file() and path.suffix.lower() == ".zip":
            zip_dest = tmp_root / path.stem
            for local, relative in extract_supported_from_zip(path, zip_dest):
                collected.append((local, f"{path.name}::{relative}"))
        elif path.is_dir():
            for candidate in iter_files([path]):
                try:
                    display = str(candidate.relative_to(path))
                except Exception:
                    display = candidate.name
                collected.append((candidate, display))
        elif path.is_file() and path.suffix.lower() in SUPPORTED_EXTENSIONS:
            collected.append((path, path.name))
    return collected


def main() -> None:
    parser = argparse.ArgumentParser(
        description=(
            "Estimate Professional Tasks source-context depth and maximum context risk. "
            "All CLI input files are treated as required files."
        )
    )
    parser.add_argument("paths", nargs="+", help="Files, folders, and/or ZIP archives to scan.")
    parser.add_argument("--prompt-file", help="Optional prompt text file.")
    parser.add_argument("--min-files", type=int, default=DEFAULT_MIN_FILES)
    parser.add_argument("--min-tokens", type=int, default=DEFAULT_MIN_SOURCE_TOKENS)
    parser.add_argument("--max-tokens", type=int, default=DEFAULT_MAX_CONTEXT_TOKENS)
    parser.add_argument("--json-out", help="Optional JSON report path.")
    args = parser.parse_args()

    if args.min_files <= 0:
        parser.error("--min-files must be greater than zero")
    if args.min_tokens <= 0:
        parser.error("--min-tokens must be greater than zero")
    if args.max_tokens <= args.min_tokens:
        parser.error("--max-tokens must be greater than --min-tokens")

    prompt_estimate: Optional[FileEstimate] = None
    if args.prompt_file:
        prompt_path = Path(args.prompt_file)
        if not prompt_path.is_file():
            parser.error(f"Prompt file not found: {prompt_path}")
        prompt_estimate = estimate_prompt(
            prompt_path.read_text(encoding="utf-8", errors="replace")
        )

    with tempfile.TemporaryDirectory(prefix="professional_tasks_context_") as tmpdir:
        collected = _collect_cli_files([Path(value) for value in args.paths], Path(tmpdir))
        estimates = [
            estimate_file(local, args.max_tokens, display_name=display)
            for local, display in collected
        ]

    if not estimates:
        print("No supported files were found.")
        return

    summary = summarize_project(
        required_estimates=estimates,
        all_estimates=estimates,
        prompt_estimate=prompt_estimate,
        min_files=args.min_files,
        min_source_tokens=args.min_tokens,
        max_context_tokens=args.max_tokens,
    )
    print_table(estimates, summary)

    if args.json_out:
        report = {
            "summary": summary,
            "files": [asdict(item) for item in estimates],
            "prompt": asdict(prompt_estimate) if prompt_estimate else None,
        }
        output_path = Path(args.json_out)
        output_path.write_text(
            json.dumps(report, indent=2, ensure_ascii=False), encoding="utf-8"
        )
        print(f"\nJSON report written to: {output_path}")


if __name__ == "__main__":
    main()
